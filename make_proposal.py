from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# --- カラーパレット ---
C_NAVY   = RGBColor(0x1A, 0x37, 0x5E)
C_GOLD   = RGBColor(0xC8, 0x9B, 0x3C)
C_LIGHT  = RGBColor(0xF4, 0xF6, 0xFA)
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_TEXT   = RGBColor(0x1A, 0x1A, 0x2E)
C_GREEN  = RGBColor(0x27, 0x7D, 0x52)
C_RED    = RGBColor(0xC0, 0x39, 0x2B)
C_BLUE2  = RGBColor(0x21, 0x6B, 0xAE)
C_GRAY   = RGBColor(0x6C, 0x75, 0x7D)
C_YELLOW = RGBColor(0xFF, 0xF3, 0xCD)

# スライドサイズ: 13.33 x 7.5 inch (ワイド16:9)
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

# ─── ユーティリティ ───────────────────────────────────────
def add_rect(slide, l, t, w, h, fill=None, line=None, line_w=None):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        if line_w:
            shape.line.width = line_w
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, l, t, w, h, size=14, bold=False, color=C_TEXT,
             align=PP_ALIGN.LEFT, italic=False):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    # 内部マージンをゼロにして有効高さを最大化
    tf.margin_top    = 0
    tf.margin_bottom = 0
    tf.margin_left   = Pt(2)
    tf.margin_right  = Pt(2)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb

def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_badge(slide, text, l, t, w=1.1, h=0.28, bg=C_NAVY, fg=C_WHITE, size=9):
    r = add_rect(slide, l, t, w, h, fill=bg)
    tf = r.text_frame
    tf.word_wrap = False
    tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = True
    run.font.color.rgb = fg

def slide_header(slide, title, subtitle=None):
    add_rect(slide, 0, 0, 13.33, 1.05, fill=C_NAVY)
    add_rect(slide, 0, 1.05, 13.33, 0.07, fill=C_GOLD)
    add_text(slide, title, 0.35, 0.1, 12.5, 0.62, size=24, bold=True, color=C_WHITE)
    if subtitle:
        add_text(slide, subtitle, 0.35, 0.68, 12.5, 0.36, size=11, color=C_GOLD)

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 1: 表紙
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
sl = prs.slides.add_slide(BLANK)
set_bg(sl, C_NAVY)

add_rect(sl, 0, 5.75, 13.33, 0.08, fill=C_GOLD)
add_rect(sl, 0, 5.9,  13.33, 1.6,  fill=RGBColor(0x10, 0x22, 0x40))

add_text(sl, "正月 親族集まり",    1.0, 1.3, 11.33, 1.1,
         size=50, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
add_text(sl, "プランニング提案書", 1.0, 2.4, 11.33, 0.9,
         size=34, bold=True, color=C_GOLD,  align=PP_ALIGN.CENTER)
add_text(sl, "2027年 1月3日（日）｜加古川・明石エリア", 1.0, 3.4, 11.33, 0.52,
         size=17, color=C_WHITE, align=PP_ALIGN.CENTER)
add_text(sl, "大人 10名 ＋ 子供 4名　計 14名",          1.0, 3.9, 11.33, 0.48,
         size=15, color=RGBColor(0xAA, 0xCC, 0xEE), align=PP_ALIGN.CENTER)
add_text(sl, "Family Planning Document  —  Confidential",
         1.0, 7.0, 11.33, 0.36, size=10, color=C_GRAY, align=PP_ALIGN.CENTER)

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 2: 基本条件
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
sl = prs.slides.add_slide(BLANK)
set_bg(sl, C_LIGHT)
slide_header(sl, "基本条件・前提", "Planning Conditions")

cards = [
    ("📅", "日程",       "2027年1月3日（仮）"),
    ("👥", "人数",       "大人10名・子供4名\n計 14名"),
    ("📍", "エリア",     "兵庫県 加古川・明石周辺"),
    ("🏠", "会場",       "自前の場所なし\n会場は別途確保が必要"),
    ("🎵", "やりたいこと", "カラオケ・トランプ等\n食事だけでなく遊びも"),
    ("🚶", "身体条件",   "足が悪いメンバーなし\n移動・階段等 問題なし"),
]

# 3列×2行のカードグリッド。余白込みで各カード幅4.1、高さ2.0
COL3 = [0.25, 4.48, 8.71]
ROW2 = [1.28, 3.42]
CW, CH = 4.1, 2.0

for i, (icon, label, body) in enumerate(cards):
    cx = COL3[i % 3]
    cy = ROW2[i // 3]
    add_rect(sl, cx, cy, CW, CH, fill=C_WHITE, line=C_GOLD, line_w=Pt(1.5))
    # ヘッダー部
    add_text(sl, icon + "  " + label,
             cx+0.14, cy+0.1, CW-0.28, 0.42,
             size=13, bold=True, color=C_NAVY)
    add_rect(sl, cx, cy+0.52, CW, 0.04, fill=C_GOLD)
    # 本文
    add_text(sl, body,
             cx+0.14, cy+0.62, CW-0.28, CH-0.72,
             size=12, color=C_TEXT)

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 3: プランA 概要
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
sl = prs.slides.add_slide(BLANK)
set_bg(sl, C_LIGHT)
slide_header(sl, "プランA　レンタルスペース ＋ テイクアウト", "Plan A – Rental Space + Takeout")

# 左パネル: 特徴リスト（x=0.25 w=4.1 y=1.28 h=6.0 = ends at 7.28）
LPX, LPY, LPW, LPH = 0.25, 1.28, 4.1, 6.0
add_rect(sl, LPX, LPY, LPW, LPH, fill=C_WHITE, line=C_BLUE2, line_w=Pt(1))
add_text(sl, "プランAの特徴", LPX+0.14, LPY+0.1, LPW-0.28, 0.4,
         size=13, bold=True, color=C_BLUE2)
add_rect(sl, LPX, LPY+0.5, LPW, 0.04, fill=C_BLUE2)

pros = [
    ("◎", "コストが最も安い",       C_GREEN),
    ("◎", "時間を自由に使える",     C_GREEN),
    ("◎", "トランプ・ゲームOK",     C_GREEN),
    ("◎", "子供が騒いでも大丈夫",   C_GREEN),
    ("◎", "飲食物の持込自由",       C_GREEN),
    ("△", "食事・飲み物の手配が必要", C_GRAY),
    ("△", "準備・片付けが発生",     C_GRAY),
    ("△", "正月は選択肢が限られる", C_GRAY),
]
for i, (mark, txt, col) in enumerate(pros):
    add_text(sl, mark + "  " + txt,
             LPX+0.14, LPY+0.65 + i*0.63, LPW-0.28, 0.52,
             size=11, color=col)

# 右パネル: 3候補カード（x=4.55 w=8.5 各カード高さ1.82 間隔1.93）
right_items = [
    {
        "name":  "A-1  ジャンカラ キャンプルーム",
        "sub":   "加古川駅 徒歩8分｜インスタベース掲載",
        "cap":   "最大 20名",
        "price": "12,100円〜 / 時間",
        "note":  "カラオケ＋ドリンクバー完備。遊びと食事が1室で完結！",
        "star":  True,
    },
    {
        "name":  "A-2  レンタルスタジオ Mi Crew",
        "sub":   "JR加古川駅 徒歩3分｜インスタベース掲載",
        "cap":   "最大 20名",
        "price": "要確認（平均 435円/人・時）",
        "note":  "駅近で集合しやすい多目的スペース。持込自由。",
        "star":  False,
    },
    {
        "name":  "A-3  明石駅前スペース",
        "sub":   "JR明石駅 徒歩5分｜スペースマーケット掲載",
        "cap":   "最大 15〜20名",
        "price": "1,732〜2,310円 / 時間",
        "note":  "子連れOK。明石側に集まるメンバーが多い場合に便利。",
        "star":  False,
    },
]

RPX = 4.55
CARD_W = 8.55
CARD_H = 1.82
CARD_GAP = 1.93   # カード開始y間隔

for i, item in enumerate(right_items):
    ty = 1.28 + i * CARD_GAP
    bg = C_YELLOW if item["star"] else C_WHITE
    ln = C_GOLD   if item["star"] else RGBColor(0xCC, 0xCC, 0xCC)
    lw = Pt(2) if item["star"] else Pt(1)
    add_rect(sl, RPX, ty, CARD_W, CARD_H, fill=bg, line=ln, line_w=lw)

    if item["star"]:
        add_badge(sl, "★ おすすめ", RPX + CARD_W - 1.3, ty + 0.1,
                  w=1.22, h=0.28, bg=C_GOLD, fg=C_WHITE, size=10)

    add_text(sl, item["name"],
             RPX+0.15, ty+0.08, CARD_W-1.5, 0.4, size=13, bold=True, color=C_NAVY)
    add_text(sl, item["sub"],
             RPX+0.15, ty+0.5, CARD_W-0.3, 0.3, size=10, color=C_GRAY, italic=True)
    add_text(sl, "収容: " + item["cap"] + "　　料金: " + item["price"],
             RPX+0.15, ty+0.82, CARD_W-0.3, 0.34, size=11, color=C_TEXT)
    add_text(sl, item["note"],
             RPX+0.15, ty+1.18, CARD_W-0.3, 0.52, size=11, color=C_TEXT)

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 4: プランA 食事・費用
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
sl = prs.slides.add_slide(BLANK)
set_bg(sl, C_LIGHT)
slide_header(sl, "プランA　食事手配 ＆ 費用概算", "Plan A – Food & Cost")

# ── 左: 食事手配 (x=0.25 w=6.1 y=1.28 h=6.0) ──
LPX, LPY, LPW, LPH = 0.25, 1.28, 6.1, 6.0
add_rect(sl, LPX, LPY, LPW, LPH, fill=C_WHITE, line=C_BLUE2, line_w=Pt(1))
add_text(sl, "食事の手配方法",
         LPX+0.14, LPY+0.1, LPW-0.28, 0.4, size=14, bold=True, color=C_BLUE2)
add_rect(sl, LPX, LPY+0.5, LPW, 0.04, fill=C_BLUE2)

# 各食事エントリー（バッジ+テキスト）
# y: 1.28+0.54 = 1.82 スタート
food_entries = [
    # (badge_text, badge_bg, title, detail)
    ("おすすめ", C_GREEN,
     "テイクアウト分担方式",
     "複数店舗に分散手配してリスク分散"),
    ("主食",    C_NAVY,
     "ごんた寿し（加古川）",
     "出前・仕出し桶 ／ 正月は要事前予約"),
    ("主食",    C_NAVY,
     "はま寿司 テイクアウト",
     "桶寿司ネット予約 ／ 大量注文向き"),
    ("副菜",    C_BLUE2,
     "善乃（加古川）",
     "オードブル・揚げ物・仕出し対応"),
    ("補助",    C_GOLD,
     "デリバリー（Uber Eats / 出前館）",
     "正月は遅延リスクあり。補填用に活用"),
    ("前日",    C_GRAY,
     "業務スーパー・コストコで調達",
     "飲み物・お菓子・紙皿・ゴミ袋等"),
]
ENTRY_H   = 0.82   # 1エントリの高さ
ENTRY_START = LPY + 0.58

for i, (badge, badge_bg, title, detail) in enumerate(food_entries):
    ey = ENTRY_START + i * ENTRY_H
    add_badge(sl, badge, LPX+0.12, ey+0.06, w=0.88, h=0.26, bg=badge_bg, fg=C_WHITE, size=8)
    add_text(sl, title,
             LPX+1.08, ey+0.04, LPW-1.22, 0.34, size=11, bold=True, color=C_TEXT)
    add_text(sl, detail,
             LPX+1.08, ey+0.38, LPW-1.22, 0.36, size=10, color=C_GRAY)
    if i < len(food_entries) - 1:
        add_rect(sl, LPX+0.12, ey+ENTRY_H-0.04, LPW-0.24, 0.02,
                 fill=RGBColor(0xDD, 0xDD, 0xDD))

# ── 右: 費用 + 持ち物 (x=6.55 w=6.55 y=1.28 h=6.0) ──
RPX, RPY, RPW, RPH = 6.55, 1.28, 6.55, 6.0
add_rect(sl, RPX, RPY, RPW, RPH, fill=C_WHITE, line=C_GOLD, line_w=Pt(1))
add_text(sl, "費用概算（14名）",
         RPX+0.18, RPY+0.1, RPW-0.36, 0.4, size=14, bold=True, color=C_NAVY)
add_rect(sl, RPX, RPY+0.5, RPW, 0.04, fill=C_GOLD)

# 費用明細 — 固定位置で配置
cost_rows = [
    ("レンタルスペース（5〜6時間）", "60,000円〜"),
    ("寿司テイクアウト × 3〜4本",   "25,000〜35,000円"),
    ("オードブル・揚げ物",           "10,000〜15,000円"),
    ("飲み物・お菓子・消耗品",       "8,000〜12,000円"),
]
CROW_START = RPY + 0.58
CROW_H     = 0.46  # 行高

for i, (label, val) in enumerate(cost_rows):
    ry = CROW_START + i * CROW_H
    add_text(sl, label,
             RPX+0.18, ry, RPW-2.5, 0.36, size=11, color=C_TEXT)
    add_text(sl, val,
             RPX+RPW-2.2, ry, 2.0, 0.36, size=11, color=C_TEXT, align=PP_ALIGN.RIGHT)
    add_rect(sl, RPX+0.12, ry+0.36, RPW-0.24, 0.02, fill=RGBColor(0xDD, 0xDD, 0xDD))

# 合計・1人あたり (固定y)
TOTAL_Y = CROW_START + len(cost_rows) * CROW_H + 0.1
add_rect(sl, RPX, TOTAL_Y, RPW, 0.55, fill=C_NAVY)
add_text(sl, "合計（推定）",
         RPX+0.18, TOTAL_Y+0.08, RPW-2.5, 0.38, size=13, bold=True, color=C_WHITE)
add_text(sl, "95,000〜112,000円",
         RPX+RPW-2.8, TOTAL_Y+0.08, 2.6, 0.38,
         size=13, bold=True, color=C_GOLD, align=PP_ALIGN.RIGHT)

PER_Y = TOTAL_Y + 0.65
add_rect(sl, RPX, PER_Y, RPW, 0.52, fill=C_YELLOW, line=C_GOLD, line_w=Pt(1))
add_text(sl, "1人あたり（概算）",
         RPX+0.18, PER_Y+0.08, RPW-2.5, 0.36, size=13, bold=True, color=C_NAVY)
add_text(sl, "約 6,800〜8,000円",
         RPX+RPW-2.8, PER_Y+0.08, 2.6, 0.36,
         size=13, bold=True, color=C_GREEN, align=PP_ALIGN.RIGHT)

# 持ち物リスト — 固定y。PER_Y+0.52+0.2 = PER_Y+0.72 スタート
MOCHI_Y = PER_Y + 0.72
add_text(sl, "当日の持ち物",
         RPX+0.18, MOCHI_Y, RPW-0.36, 0.34, size=12, bold=True, color=C_NAVY)
add_rect(sl, RPX+0.12, MOCHI_Y+0.34, RPW-0.24, 0.03, fill=C_GOLD)

mochi_items = [
    "紙皿・紙コップ・割り箸（多めに）",
    "ゴミ袋（大）× 数枚",
    "トランプ・UNO・ボードゲーム",
    "Bluetoothスピーカー（BGM用）",
    "延長コード・ウェットティッシュ",
]
MOCHI_ITEM_H = 0.32  # 1アイテムの行高（5項目が panel 内に収まるよう調整）

for j, it in enumerate(mochi_items):
    iy = MOCHI_Y + 0.44 + j * MOCHI_ITEM_H
    add_text(sl, "✔  " + it,
             RPX+0.18, iy, RPW-0.36, 0.3, size=10, color=C_TEXT)

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 5: プランB
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
sl = prs.slides.add_slide(BLANK)
set_bg(sl, C_LIGHT)
slide_header(sl, "プランB　ホテル宴会場（加古川プラザホテル）", "Plan B – Hotel Banquet Hall")

# ── 左: ホテル情報 (x=0.25 w=8.1 y=1.28 h=3.45) ──
add_rect(sl, 0.25, 1.28, 8.1, 3.45, fill=C_WHITE, line=C_NAVY, line_w=Pt(1.5))
add_text(sl, "加古川プラザホテル",
         0.4, 1.38, 7.8, 0.5, size=20, bold=True, color=C_NAVY)
add_rect(sl, 0.25, 1.88, 8.1, 0.04, fill=C_GOLD)

hotel_info = [
    ("住所",     "加古川市加古川町溝之口800"),
    ("宴会 TEL", "079-421-6012　（受付 10:00〜17:00）"),
    ("宴会場",   "大・中・小（芙蓉の間・牡丹の間 等）"),
    ("雰囲気",   "シャンデリアあり ／ フォーマル〜セミフォーマル"),
    ("駐車場",   "あり（台数は要確認）"),
]
for i, (k, v) in enumerate(hotel_info):
    hy = 2.01 + i * 0.52
    add_text(sl, k, 0.4, hy, 1.75, 0.4, size=11, bold=True, color=C_GRAY)
    add_text(sl, v, 2.2, hy, 6.0,  0.4, size=11, color=C_TEXT)

# ── 右: メリット・デメリット (x=8.55 w=4.55 y=1.28 h=3.45) ──
add_rect(sl, 8.55, 1.28, 4.55, 3.45, fill=C_WHITE,
         line=RGBColor(0xCC, 0xCC, 0xCC), line_w=Pt(1))
add_text(sl, "メリット・デメリット",
         8.7, 1.38, 4.25, 0.4, size=13, bold=True, color=C_NAVY)
add_rect(sl, 8.55, 1.88, 4.55, 0.04, fill=C_GOLD)

bds = [
    ("◎", "準備・片付けが不要",     C_GREEN),
    ("◎", "料理の質・特別感が高い", C_GREEN),
    ("◎", "正月宴会プランあり",     C_GREEN),
    ("△", "カラオケは別途手配",     C_GRAY),
    ("△", "滞在時間は2〜3時間が目安", C_GRAY),
    ("△", "費用がやや高い",         C_GRAY),
]
for i, (ico, txt, col) in enumerate(bds):
    add_text(sl, ico + "  " + txt,
             8.7, 2.01 + i * 0.52, 4.25, 0.44, size=11, color=col)

# ── 費用表 (y=4.9 h=2.38) ──
TBX, TBY, TBW, TBH = 0.25, 4.9, 12.85, 2.38
add_rect(sl, TBX, TBY, TBW, TBH, fill=C_WHITE, line=C_GOLD, line_w=Pt(1))
add_text(sl, "費用概算（14名）",
         TBX+0.18, TBY+0.1, TBW-0.36, 0.38, size=14, bold=True, color=C_NAVY)
add_rect(sl, TBX, TBY+0.5, TBW, 0.04, fill=C_GOLD)

# テーブルヘッダー
plan_col_x = [TBX+0.18, TBX+5.5, TBX+9.5]
plan_col_w = [5.1, 3.8, 3.1]
plan_hdr   = ["プランタイプ", "1名あたり", "14名合計"]
for j, (hd, cx, cw) in enumerate(zip(plan_hdr, plan_col_x, plan_col_w)):
    bg = C_NAVY if j == 0 else (C_BLUE2 if j == 1 else C_GREEN)
    add_rect(sl, cx, TBY+0.56, cw, 0.36, fill=bg)
    add_text(sl, hd, cx+0.05, TBY+0.6, cw-0.1, 0.3,
             size=11, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

plan_rows = [
    ("和食コース（料理・飲物込み）", "5,000〜8,000円", "70,000〜112,000円", False),
    ("バイキング形式",               "3,000〜5,000円", "42,000〜70,000円",  True),
    ("会場のみ（持込）",             "1,000〜2,000円", "14,000〜28,000円",  False),
]
for i, (a, b, c, highlight) in enumerate(plan_rows):
    ry = TBY + 0.92 + i * 0.44
    bg = C_YELLOW if highlight else (C_WHITE if i % 2 == 0 else RGBColor(0xF0, 0xF4, 0xF9))
    add_rect(sl, TBX, ry-0.04, TBW, 0.44, fill=bg)
    add_text(sl, a, plan_col_x[0]+0.05, ry, plan_col_w[0]-0.1, 0.36, size=11, color=C_TEXT)
    add_text(sl, b, plan_col_x[1]+0.05, ry, plan_col_w[1]-0.1, 0.36,
             size=11, color=C_TEXT, align=PP_ALIGN.CENTER)
    fc = C_GREEN if highlight else C_TEXT
    add_text(sl, c, plan_col_x[2]+0.05, ry, plan_col_w[2]-0.1, 0.36,
             size=11, bold=highlight, color=fc, align=PP_ALIGN.CENTER)

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 6: 比較表
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
sl = prs.slides.add_slide(BLANK)
set_bg(sl, C_LIGHT)
slide_header(sl, "プランA vs プランB　比較", "Comparison")

cmp_col_x = [0.25, 3.85, 9.05]
cmp_col_w = [3.45, 5.05, 4.05]
cmp_hdr   = ["比較軸", "A  レンタルスペース", "B  ホテル宴会場"]

for j, (hd, cx, cw) in enumerate(zip(cmp_hdr, cmp_col_x, cmp_col_w)):
    bg = C_NAVY if j == 0 else (C_BLUE2 if j == 1 else C_GREEN)
    add_rect(sl, cx, 1.28, cw, 0.48, fill=bg)
    add_text(sl, hd, cx+0.1, 1.32, cw-0.2, 0.38,
             size=13, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

cmp_rows = [
    ("費用",          "◎ 安い（1人 6,800〜8,000円）",  "△ やや高い（1人 3,000〜8,000円）"),
    ("準備・片付け",  "△ 自分たちで対応",              "◎ ほぼ不要"),
    ("カラオケ・遊び","◎ A-1なら同室で可能",            "✕ 別途移動が必要"),
    ("子供の自由度",  "◎ 騒いでもOK",                  "△ 場所による"),
    ("長居（5時間+）","◎ 時間制限なし",                 "△ 2〜3時間が目安"),
    ("料理の質",      "△ テイクアウト依存",             "◎ ホテル料理"),
    ("特別感",        "○ 自分たちらしい場",             "◎ 正月らしい格式"),
    ("予約しやすさ",  "○ オンライン即予約可",           "△ 電話 & 早期確認が必要"),
]
ROW_H = 0.54
for i, row in enumerate(cmp_rows):
    ry = 1.76 + i * ROW_H
    bg_r = C_WHITE if i % 2 == 0 else RGBColor(0xEB, 0xEF, 0xF5)
    # 行全体の背景
    add_rect(sl, 0.25, ry, sum(cmp_col_w), ROW_H, fill=bg_r)
    for j, (val, cx, cw) in enumerate(zip(row, cmp_col_x, cmp_col_w)):
        if j == 0:
            col = C_NAVY
            al  = PP_ALIGN.LEFT
            bd  = True
        else:
            col = (C_GREEN if "◎" in val else
                   C_RED   if "✕" in val else C_TEXT)
            al  = PP_ALIGN.CENTER
            bd  = False
        add_text(sl, val, cx+0.1, ry+0.08, cw-0.2, ROW_H-0.1,
                 size=11, bold=bd, color=col, align=al)

# 結論バナー（2行に分けて確実に収める）
BANNER_Y = 1.76 + len(cmp_rows) * ROW_H + 0.08
add_rect(sl, 0.25, BANNER_Y, 13.05, 0.56, fill=C_NAVY)
add_text(sl, "コスト・自由度を重視 →  プランA（A-1 ジャンカラ キャンプルーム）",
         0.4, BANNER_Y+0.02, 6.4, 0.48,
         size=11, bold=True, color=C_GOLD)
add_text(sl, "手間なし・特別感を重視 →  プランB（加古川プラザホテル）",
         6.9, BANNER_Y+0.02, 6.3, 0.48,
         size=11, bold=True, color=RGBColor(0x90, 0xD4, 0xB0))

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 7: 次のアクション
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
sl = prs.slides.add_slide(BLANK)
set_bg(sl, C_LIGHT)
slide_header(sl, "今すぐやること　Next Actions", "Action Plan")

steps = [
    {
        "no":    "01",
        "title": "加古川プラザホテルへ電話",
        "body":  "TEL: 079-421-6012（受付 10:00〜17:00）\n1月3日・14名・小宴会場の空きと料金プランを確認",
        "dead":  "今週中",
        "urgent": True,
    },
    {
        "no":    "02",
        "title": "ジャンカラ キャンプルーム 空き確認",
        "body":  "インスタベースで「加古川 キャンプルーム」を検索\n1月3日の空き状況・料金・予約方法を確認する",
        "dead":  "今週中",
        "urgent": True,
    },
    {
        "no":    "03",
        "title": "プランを1本に絞って即予約",
        "body":  "AまたはBどちらかに決定し即予約\n正月3日は人気日・早い者勝ち",
        "dead":  "今月中",
        "urgent": False,
    },
    {
        "no":    "04",
        "title": "親族へ日程案内・出欠確認",
        "body":  "LINEグループ等で確定人数を確認\n子供の年齢・食べ物の好みも把握しておく",
        "dead":  "今月中",
        "urgent": False,
    },
    {
        "no":    "05",
        "title": "食事の手配（プランA の場合）",
        "body":  "ごんた寿し・はま寿司のテイクアウトを予約\n正月テイクアウトは 11〜12月に早めに予約",
        "dead":  "11〜12月",
        "urgent": False,
    },
]

# カード配置: 2列×3行（最後の5枚目は中央）
# 行間=1.82、カード高さ=1.70、footer y=6.85
CARD_W2 = 6.15
CARD_H2 = 1.70
COL_X2  = [0.25, 6.55]
ROW_Y2  = [1.28, 3.1, 4.92]

def card_pos(idx):
    if idx == 4:
        return (13.33 - CARD_W2) / 2, ROW_Y2[2]
    return COL_X2[idx % 2], ROW_Y2[idx // 2]

for i, s in enumerate(steps):
    lx, ty = card_pos(i)
    ln_col = C_GOLD if s["urgent"] else RGBColor(0xCC, 0xCC, 0xCC)
    add_rect(sl, lx, ty, CARD_W2, CARD_H2, fill=C_WHITE, line=ln_col,
             line_w=Pt(2 if s["urgent"] else 1))

    # NO バッジ
    add_rect(sl, lx, ty, 0.48, 0.48, fill=C_NAVY)
    add_text(sl, s["no"], lx, ty+0.06, 0.48, 0.36,
             size=12, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    # 期限バッジ
    badge_bg = C_RED if s["urgent"] else C_BLUE2
    add_badge(sl, s["dead"], lx + CARD_W2 - 1.28, ty + 0.1,
              w=1.2, h=0.28, bg=badge_bg, fg=C_WHITE, size=9)

    # タイトル
    add_text(sl, s["title"], lx+0.55, ty+0.07, CARD_W2-1.85, 0.38,
             size=12, bold=True, color=C_NAVY)

    # 本文（2行、フォント10pt×0.22inch/行=0.44inch。h=0.9で余裕あり）
    add_text(sl, s["body"], lx+0.12, ty+0.52, CARD_W2-0.24, 0.9,
             size=10, color=C_TEXT)

# フッターバー (y=6.86 h=0.46、カード末尾は 4.92+1.70=6.62 → 重ならない)
FOOTER_Y = 6.86
add_rect(sl, 0.25, FOOTER_Y, 13.05, 0.46, fill=C_NAVY)
add_text(sl,
         "⚠  1月3日は正月繁忙期。場所確保は今すぐ。決まったら食事手配に即移行する。",
         0.45, FOOTER_Y+0.06, 12.85, 0.34,
         size=11, bold=True, color=C_GOLD, align=PP_ALIGN.CENTER)

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# 保存
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
out = "/home/user/test/正月集まり_プランニング提案書.pptx"
prs.save(out)
print("Done:", out)
