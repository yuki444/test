"""PowerPointプレゼンテーション作成ユーティリティ"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import json
import sys


def create_presentation(slides_data: list[dict], output_path: str = "output.pptx") -> str:
    """
    スライドデータからPowerPointファイルを生成する。

    slides_data: [
        {
            "title": "スライドタイトル",
            "content": ["箇条書き1", "箇条書き2"],
            "notes": "発表者ノート（任意）"
        },
        ...
    ]
    """
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    title_slide_layout = prs.slide_layouts[0]
    content_slide_layout = prs.slide_layouts[1]

    for i, slide_data in enumerate(slides_data):
        layout = title_slide_layout if i == 0 else content_slide_layout
        slide = prs.slides.add_slide(layout)

        title = slide.shapes.title
        if title:
            title.text = slide_data.get("title", "")
            title.text_frame.paragraphs[0].font.bold = True

        if i == 0:
            if len(slide.placeholders) > 1:
                subtitle = slide.placeholders[1]
                subtitle.text = slide_data.get("subtitle", "")
        else:
            if len(slide.placeholders) > 1:
                body = slide.placeholders[1]
                tf = body.text_frame
                tf.clear()
                content = slide_data.get("content", [])
                for j, line in enumerate(content):
                    if j == 0:
                        tf.text = line
                    else:
                        p = tf.add_paragraph()
                        p.text = line
                    tf.paragraphs[j].level = slide_data.get("levels", [0] * len(content))[j] if "levels" in slide_data else 0

        if "notes" in slide_data and slide_data["notes"]:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = slide_data["notes"]

    prs.save(output_path)
    return output_path


def main():
    if len(sys.argv) < 2:
        # サンプル実行
        sample_slides = [
            {
                "title": "プレゼンテーションタイトル",
                "subtitle": "サブタイトルをここに入力"
            },
            {
                "title": "アジェンダ",
                "content": ["はじめに", "本題", "まとめ"]
            },
            {
                "title": "本題",
                "content": [
                    "ポイント1: 重要な内容",
                    "ポイント2: 詳細説明",
                    "ポイント3: 結論"
                ],
                "notes": "ここに発表者ノートを入力"
            },
            {
                "title": "まとめ",
                "content": ["今日学んだこと", "次のステップ", "ご清聴ありがとうございました"]
            }
        ]
        output = create_presentation(sample_slides, "sample.pptx")
        print(f"サンプルプレゼンテーションを作成しました: {output}")
    else:
        json_path = sys.argv[1]
        output_path = sys.argv[2] if len(sys.argv) > 2 else "output.pptx"
        with open(json_path, "r", encoding="utf-8") as f:
            slides_data = json.load(f)
        output = create_presentation(slides_data, output_path)
        print(f"プレゼンテーションを作成しました: {output}")


if __name__ == "__main__":
    main()
